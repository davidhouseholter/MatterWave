#!/usr/bin/env python3
"""
Simple Markdown-to-PPTX exporter for the Deck.md in this repo.

This script is intentionally minimal and converts each Reveal-style slide
separated by a line with exactly '---' into one PPTX slide. It extracts a
title from the first Markdown heading in the slide (a line starting with
'# ') and converts list items (lines starting with '- ' or '* ') into bullets.

Requirements: python-pptx

Usage:
  python export_to_pptx.py Deck.md Deck.pptx

Notes:
- This is not a full Markdown renderer; it handles headings, paragraphs, and
  simple bulleted lists. It's designed to produce a presentation suitable for
  editing and distribution when Pandoc is not available.
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def split_slides(md: str):
    # Remove YAML frontmatter if present
    if md.startswith("---"):
        # find the second '---' delimiter
        parts = md.split('\n')
        # find index of second '---'
        try:
            second = parts.index('---', 1)
            md = '\n'.join(parts[second+1:])
        except ValueError:
            # no second delimiter; leave as-is
            pass

    slides = []
    current = []
    for line in md.splitlines():
        if line.strip() == '---':
            if current:
                slides.append('\n'.join(current).strip())
                current = []
            continue
        current.append(line)
    if current:
        slides.append('\n'.join(current).strip())
    return slides


def extract_title_and_bullets(slide_md: str):
    lines = [l.rstrip() for l in slide_md.splitlines() if l.strip() != '']
    title = None
    bullets = []
    body_lines = []
    for i, l in enumerate(lines):
        if l.startswith('# '):
            title = l[2:].strip()
            body_lines = lines[i+1:]
            break
    if title is None:
        # fallback: use first non-empty line as title
        if lines:
            title = lines[0]
            body_lines = lines[1:]
        else:
            title = ''

    for l in body_lines:
        stripped = l.lstrip()
        if stripped.startswith('- ') or stripped.startswith('* '):
            bullets.append(stripped[2:].strip())
        else:
            # treat as paragraph; break into sentences for bullets if long
            bullets.append(stripped)

    return title, bullets


def make_pptx(deck_md_path: Path, out_pptx_path: Path):
    md = deck_md_path.read_text(encoding='utf-8')
    slides = split_slides(md)

    prs = Presentation()
    # ensure at least one slide layout exists
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    first = True
    for s in slides:
        title, bullets = extract_title_and_bullets(s)
        if first:
            slide = prs.slides.add_slide(title_layout)
            if title:
                slide.shapes.title.text = title
            # add a small subtitle with project name if available
            try:
                subtitle = slide.placeholders[1]
                subtitle.text = 'MatterWave — KG-CAE'
            except Exception:
                pass
            first = False
        else:
            slide = prs.slides.add_slide(body_layout)
            if title:
                slide.shapes.title.text = title
            # add bullets
            try:
                body = slide.shapes.placeholders[1].text_frame
            except Exception:
                # find a textbox to place bullets
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
                body = txBox.text_frame
            body.clear()
            for b in bullets:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(18)

    prs.save(str(out_pptx_path))


def main(argv):
    if len(argv) < 3:
        print('Usage: export_to_pptx.py Deck.md Deck.pptx')
        return 1
    deck = Path(argv[1])
    out = Path(argv[2])
    if not deck.exists():
        print('Deck file not found:', deck)
        return 2
    make_pptx(deck, out)
    print('Wrote', out)
    return 0


if __name__ == '__main__':
    raise SystemExit(main(sys.argv))
